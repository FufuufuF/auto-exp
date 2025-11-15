from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table

from pypdf import PdfReader

from pptx import Presentation

class MicrosoftOfficeReader:
    
    @staticmethod
    def read_microsoft_word(file_path: str) -> str:
        doc = Document(file_path)
        full_text = ""
        for block in doc.iter_inner_content():
            if isinstance(block, Paragraph):
                text = block.text.strip()
                if text:
                    full_text += text + '\n'
            elif isinstance(block, Table):
                for row in block.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            full_text += text + '\t'
                    full_text += '\n'
        return full_text

    @staticmethod
    def read_microsoft_pdf(file_path: str) -> str:
        try:
            reader = PdfReader(file_path)
            full_text = ""
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + '\n'
                else:
                    print(f"第{page.page_number}页没有文本")
            return full_text
        except Exception as e:
            print(f"读取PDF失败: {e}")
            return f"读取PDF失败: {e}"

    @staticmethod
    def read_microsoft_ppt(file_path: str) -> str:
        try:
            presentation = Presentation(file_path)
            full_text = []
            for slide_num, slide in enumerate(presentation.slides):
                devide_str = f"\n--- 第 {slide_num + 1} 张幻灯片 ---"
                slide_content = [devide_str]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                    
                    if shape.has_table:
                        table_content = []
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text_frame.text.strip() for cell in row.cells]
                            table_content.append(" | ".join(row_text))
                        slide_content.append("\n[表格内容]:\n" + "\n".join(table_content))

                full_text.append('\n'.join(slide_content))
            return '\n'.join(full_text)
            
        except Exception as e:
            print(f"读取PPT失败: {e}")
            return f"读取PPT失败: {e}"

if __name__ == "__main__":
    file_path = "test_ppt.pptx"
    full_text = MicrosoftOfficeReader.read_microsoft_ppt(file_path)
    print(full_text)

